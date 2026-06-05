from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from zipfile import ZipFile

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader


@dataclass(frozen=True)
class ExtractedSection:
    label: str
    text: str


def extract_sections(file_path: Path, extension: str) -> list[ExtractedSection]:
    if extension == ".pdf":
        return _extract_pdf(file_path)
    if extension == ".docx":
        return _extract_docx(file_path)
    if extension == ".pptx":
        return _extract_pptx(file_path)
    raise ValueError("Unsupported file type.")


def chunk_sections(sections: list[ExtractedSection], max_chars: int = 1200) -> list[ExtractedSection]:
    chunks: list[ExtractedSection] = []
    for section in sections:
        text = " ".join(section.text.split())
        if not text:
            continue
        start = 0
        while start < len(text):
            end = min(start + max_chars, len(text))
            if end < len(text):
                split_at = text.rfind(" ", start, end)
                if split_at > start + 300:
                    end = split_at
            chunks.append(ExtractedSection(label=section.label, text=text[start:end].strip()))
            start = end
    return chunks


def _extract_pdf(file_path: Path) -> list[ExtractedSection]:
    reader = PdfReader(str(file_path))
    sections: list[ExtractedSection] = []
    for index, page in enumerate(reader.pages, start=1):
        sections.append(ExtractedSection(label=f"page {index}", text=page.extract_text() or ""))
    return sections


def _extract_docx(file_path: Path) -> list[ExtractedSection]:
    document = DocxDocument(str(file_path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return [ExtractedSection(label="document", text="\n".join(paragraphs))]


def _extract_pptx(file_path: Path) -> list[ExtractedSection]:
    presentation = Presentation(str(file_path))
    sections: list[ExtractedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        sections.append(ExtractedSection(label=f"slide {index}", text="\n".join(texts)))
    return sections


def validate_office_file(file_path: Path) -> None:
    # Forces a quick ZIP read so corrupt Office files fail before parser-specific code.
    with ZipFile(file_path):
        return
