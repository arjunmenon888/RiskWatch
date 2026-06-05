from __future__ import annotations

import shutil
import tempfile
import time
from pathlib import Path

import httpx
from docx import Document as DocxDocument
from pptx import Presentation
from sqlalchemy import delete

from app.db.session import SessionLocal
from app.models.document import Document
from app.models.game import Game
from app.models.user import User

BASE_URL = "http://127.0.0.1:8002"


def main() -> None:
    stamp = int(time.time())
    temp_dir = Path(tempfile.mkdtemp(dir="C:/tmp"))
    files = create_test_files(temp_dir)
    email = f"phase3{stamp}@example.com"
    game_id: int | None = None

    try:
        with httpx.Client(timeout=20) as client:
            signup = client.post(
                f"{BASE_URL}/auth/signup",
                json={
                    "email": email,
                    "full_name": "Phase Three",
                    "password": "password123",
                    "role": "creator",
                },
            )
            signup.raise_for_status()
            token = signup.json()["access_token"]
            headers = {"Authorization": f"Bearer {token}"}

            game = client.post(
                f"{BASE_URL}/creator/games",
                headers=headers,
                json={
                    "title": "Phase Three Upload Game",
                    "description": "upload smoke",
                    "category": "Safety",
                    "visibility": "private",
                    "creation_mode": "ai",
                },
            )
            game.raise_for_status()
            game_id = game.json()["id"]

            results = []
            for path, mime_type in files:
                with path.open("rb") as file:
                    upload = client.post(
                        f"{BASE_URL}/api/games/{game_id}/documents/upload",
                        headers=headers,
                        files={"file": (path.name, file, mime_type)},
                    )
                upload.raise_for_status()
                body = upload.json()
                results.append(
                    (
                        path.suffix,
                        upload.status_code,
                        body["status"],
                        body["chunk_count"],
                        body["extracted_text_char_count"],
                    )
                )

            documents = client.get(f"{BASE_URL}/api/games/{game_id}/documents", headers=headers)
            documents.raise_for_status()
            print("game", game_id)
            print("uploads", results)
            print("list_documents", len(documents.json()))
    finally:
        cleanup(email, game_id)
        shutil.rmtree(temp_dir, ignore_errors=True)


def create_test_files(temp_dir: Path) -> list[tuple[Path, str]]:
    pdf = temp_dir / "phase3.pdf"
    docx = temp_dir / "phase3.docx"
    pptx = temp_dir / "phase3.pptx"

    write_simple_pdf(pdf, "Phase three PDF safety content for extraction and chunk creation.")

    document = DocxDocument()
    document.add_paragraph("Phase three DOCX safety content for extraction and chunks.")
    document.save(docx)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Phase three PPTX"
    textbox = slide.shapes.add_textbox(1_000_000, 1_500_000, 6_000_000, 1_000_000)
    textbox.text_frame.text = "PPTX slide safety content for extraction and chunks."
    presentation.save(pptx)

    return [
        (pdf, "application/pdf"),
        (docx, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        (pptx, "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ]


def write_simple_pdf(path: Path, text: str) -> None:
    stream = f"BT /F1 18 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    data = b"%PDF-1.4\n"
    offsets = []
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(data))
        data += f"{index} 0 obj\n".encode() + obj + b"\nendobj\n"

    xref = len(data)
    data += f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode()
    for offset in offsets:
        data += f"{offset:010d} 00000 n \n".encode()
    data += f"trailer << /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    path.write_bytes(data)


def cleanup(email: str, game_id: int | None) -> None:
    with SessionLocal() as db:
        if game_id is not None:
            db.execute(delete(Document).where(Document.game_id == game_id))
            db.execute(delete(Game).where(Game.id == game_id))
            shutil.rmtree(Path("storage") / "games" / str(game_id), ignore_errors=True)
        db.execute(delete(User).where(User.email == email))
        db.commit()


if __name__ == "__main__":
    main()
